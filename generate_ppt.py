"""
基于 Java 的手写体收集和管理系统 - 答辩 PPT 生成脚本
技术栈：python-pptx
作者：handwriting-team
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============== 视觉规范 ==============
PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)        # #1E3A5F 深蓝
ACCENT = RGBColor(0xC8, 0x9B, 0x3C)         # #C89B3C 暖金
TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)
TEXT_GRAY = RGBColor(0x6B, 0x72, 0x80)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFB)
BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
GOLD_SOFT = RGBColor(0xFE, 0xF3, 0xC7)

CN_FONT = "微软雅黑"
EN_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============== 工具函数 ==============
def set_font(run, name=CN_FONT, size=18, bold=False, color=TEXT_DARK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    eastAsia = rPr.find(qn('a:ea'))
    if eastAsia is None:
        eastAsia = etree.SubElement(rPr, qn('a:ea'))
    eastAsia.set('typeface', name)


def add_textbox(slide, x, y, w, h, text, size=18, bold=False,
                color=TEXT_DARK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font=CN_FONT, line_spacing=1.3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        set_font(run, name=font, size=size, bold=bold, color=color)
    return tb


def add_rect(slide, x, y, w, h, fill, line=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width:
            shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def add_rounded_rect(slide, x, y, w, h, fill, line=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width:
            shape.line.width = line_width
    shape.shadow.inherit = False
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_header(slide, page_no, total_pages, section_title, section_no):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.05), ACCENT)
    add_textbox(
        slide, Inches(0.5), Inches(0.18), Inches(8), Inches(0.35),
        f"PART {section_no:02d}  ·  {section_title}" if section_no > 0 else section_title,
        size=11, bold=True, color=ACCENT, font=EN_FONT
    )
    add_textbox(
        slide, Inches(9), Inches(0.18), Inches(4), Inches(0.35),
        f"HandWrite System  |  {page_no} / {total_pages}",
        size=10, color=TEXT_GRAY, align=PP_ALIGN.RIGHT, font=EN_FONT
    )


def add_slide_title(slide, title, subtitle=None):
    add_textbox(
        slide, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.7),
        title, size=32, bold=True, color=PRIMARY
    )
    add_rect(slide, Inches(0.5), Inches(1.45), Inches(0.6), Inches(0.06), ACCENT)
    if subtitle:
        add_textbox(
            slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.4),
            subtitle, size=14, color=TEXT_GRAY
        )


# ============== 封面 ==============
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)
    add_rect(slide, Inches(11.5), 0, Inches(0.04), SLIDE_H, ACCENT)
    add_rect(slide, 0, Inches(7.3), SLIDE_W, Inches(0.04), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.7), Inches(4), Inches(0.4),
                "COURSE DESIGN DEFENSE", size=12, bold=True, color=ACCENT, font=EN_FONT)
    add_textbox(slide, Inches(0.8), Inches(1.05), Inches(4), Inches(0.3),
                "课程设计答辩  ·  2026.06.22", size=11, color=LIGHT_GRAY)

    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2),
                "基于 Java 的手写体收集和管理系统",
                size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.6),
                "A Handwriting Collection & Management Platform",
                size=20, color=ACCENT, font=EN_FONT)
    add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.5),
                "Spring Boot 3  ·  Vue 3  ·  MyBatis-Plus  ·  MinIO  ·  JWT",
                size=14, color=LIGHT_GRAY, font=EN_FONT)

    add_rect(slide, Inches(0.8), Inches(4.7), Inches(2.5), Inches(0.04), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(5.0), Inches(8), Inches(0.4),
                "答辩人：XXX", size=16, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(8), Inches(0.4),
                "指导教师：XXX  教授", size=16, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(6.0), Inches(8), Inches(0.4),
                "学  院：计算机科学与技术学院", size=14, color=LIGHT_GRAY)

    add_textbox(slide, Inches(9.5), Inches(5.5), Inches(3.5), Inches(1.5),
                "墨韵", size=72, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, font=CN_FONT)


# ============== 目录 ==============
def slide_toc(prs, page_no, total_pages):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_LIGHT)
    add_header(slide, page_no, total_pages, "目录", 0)
    add_slide_title(slide, "CONTENTS", "本次答辩的内容导览")

    items = [
        ("01", "项目背景与目标", "Background & Goals"),
        ("02", "技术选型", "Technology Stack"),
        ("03", "系统架构设计", "System Architecture"),
        ("04", "关键技术实现", "Key Implementation"),
        ("05", "工程化与部署", "Engineering & Deployment"),
        ("06", "测试与质量保证", "Testing & QA"),
        ("07", "亮点总结与展望", "Highlights & Outlook"),
    ]
    for i, (no, cn, en) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(2.4) + row * Inches(0.8)
        add_textbox(slide, x, y, Inches(0.7), Inches(0.6),
                    no, size=28, bold=True, color=ACCENT, font=EN_FONT)
        add_textbox(slide, x + Inches(0.7), y + Inches(0.05), Inches(3.2), Inches(0.4),
                    cn, size=16, bold=True, color=PRIMARY)
        add_textbox(slide, x + Inches(0.7), y + Inches(0.4), Inches(3.2), Inches(0.3),
                    en, size=10, color=TEXT_GRAY, font=EN_FONT)
        if col == 0:
            add_rect(slide, x + Inches(4.5), y + Inches(0.1),
                     Inches(0.02), Inches(0.5), DIVIDER)


# ============== 章节封面 ==============
def slide_section(prs, page_no, total_pages, no, title, en):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)
    add_rect(slide, Inches(0.5), Inches(2.8), Inches(0.1), Inches(2.0), ACCENT)

    add_textbox(slide, Inches(0.9), Inches(2.6), Inches(4), Inches(0.6),
                f"PART {no:02d}", size=20, bold=True, color=ACCENT, font=EN_FONT)
    add_textbox(slide, Inches(0.9), Inches(3.2), Inches(11), Inches(1.0),
                title, size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.9), Inches(4.2), Inches(11), Inches(0.5),
                en, size=18, color=LIGHT_GRAY, font=EN_FONT)


# ============== 通用内容页 ==============
def slide_content_base(prs, page_no, total_pages, section_no, section_title, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_LIGHT)
    add_header(slide, page_no, total_pages, section_title, section_no)
    add_slide_title(slide, title, subtitle)
    return slide


# ============== 项目背景 ==============
def slide_background(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "项目背景与目标", "Background & Goals"
    )
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(0.4),
                "项目定位", size=18, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.5), Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    add_rounded_rect(slide, Inches(0.5), Inches(2.85), Inches(6.0), Inches(1.7), BG_CARD)
    add_textbox(
        slide, Inches(0.7), Inches(3.0), Inches(5.6), Inches(1.4),
        "面向中文手写体研究的样本采集与审核管理平台，\n"
        "为字体设计师、AI 手写识别研究者提供高质量的\n"
        "「在线书写 → 元数据采集 → 审核归档 → 数据可视化」\n"
        "完整工程化解决方案。",
        size=14, color=TEXT_DARK, line_spacing=1.5
    )

    add_textbox(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(0.4),
                "核心目标", size=18, bold=True, color=PRIMARY)
    add_rect(slide, Inches(6.8), Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    goals = [
        ("01", "用户体系", "JWT + RBAC 注册登录与权限隔离"),
        ("02", "样本采集", "Canvas 手写板 + 文件直传"),
        ("03", "审核工作流", "通过 / 驳回 + 审计日志"),
        ("04", "数据可视化", "趋势 / 状态 / 字典进度 / 贡献排行"),
    ]
    for i, (no, t, d) in enumerate(goals):
        y = Inches(2.85) + i * Inches(0.7)
        add_rounded_rect(slide, Inches(6.8), y, Inches(6.0), Inches(0.6), BG_CARD)
        add_textbox(slide, Inches(7.0), y + Inches(0.13), Inches(0.6), Inches(0.4),
                    no, size=16, bold=True, color=ACCENT, font=EN_FONT)
        add_textbox(slide, Inches(7.6), y + Inches(0.05), Inches(2), Inches(0.3),
                    t, size=14, bold=True, color=PRIMARY)
        add_textbox(slide, Inches(9.6), y + Inches(0.1), Inches(3.2), Inches(0.4),
                    d, size=12, color=TEXT_GRAY)

    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
                "四大核心业务场景", size=18, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.5), Inches(5.4), Inches(0.4), Inches(0.04), ACCENT)

    scenes = [
        ("📝", "在线书写", "Canvas + 压感笔锋还原"),
        ("📚", "字典管理", "字符字典增删改 + 批量导入"),
        ("✅", "审核归档", "审核员工作流 + 审计日志"),
        ("📊", "数据洞察", "ECharts 多维可视化大屏"),
    ]
    for i, (icon, t, d) in enumerate(scenes):
        x = Inches(0.5) + i * Inches(3.1)
        add_rounded_rect(slide, x, Inches(5.6), Inches(2.9), Inches(1.5), BG_CARD)
        add_textbox(slide, x + Inches(0.2), Inches(5.75), Inches(2.6), Inches(0.5),
                    icon, size=28)
        add_textbox(slide, x + Inches(0.2), Inches(6.25), Inches(2.6), Inches(0.35),
                    t, size=15, bold=True, color=PRIMARY)
        add_textbox(slide, x + Inches(0.2), Inches(6.6), Inches(2.6), Inches(0.4),
                    d, size=11, color=TEXT_GRAY)


# ============== 后端技术栈 ==============
def slide_tech_backend(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "技术选型 · 后端", "Backend Tech Stack"
    )

    card_x, card_y, card_w, card_h = Inches(0.5), Inches(2.3), Inches(6.0), Inches(4.7)
    add_rounded_rect(slide, card_x, card_y, card_w, card_h, BG_CARD)
    add_textbox(slide, card_x + Inches(0.3), card_y + Inches(0.2), Inches(5), Inches(0.4),
                "核心框架", size=14, bold=True, color=ACCENT)
    add_rect(slide, card_x + Inches(0.3), card_y + Inches(0.65), Inches(0.3), Inches(0.03), ACCENT)

    items = [
        ("JDK", "Eclipse Temurin 17 (LTS)"),
        ("Framework", "Spring Boot 3.2.5"),
        ("Persistence", "Spring Data JPA + MyBatis-Plus 3.5.5"),
        ("Database", "MySQL 8.0+"),
        ("Cache", "Spring Data Redis (Lettuce) 7.x"),
        ("Security", "Spring Security 6 + JJWT 0.12.5"),
        ("Object Storage", "MinIO 8.5 / Azure Blob"),
        ("API Doc", "springdoc-openapi 2.5.0"),
        ("DB Migration", "Flyway 9.22.3"),
        ("Utilities", "Lombok · MapStruct · Hutool"),
        ("Monitor", "Actuator + Micrometer Prometheus"),
    ]
    for i, (k, v) in enumerate(items):
        y = card_y + Inches(0.85) + i * Inches(0.32)
        add_textbox(slide, card_x + Inches(0.3), y, Inches(1.6), Inches(0.3),
                    k, size=11, color=TEXT_GRAY, font=EN_FONT)
        add_textbox(slide, card_x + Inches(1.9), y, Inches(4.0), Inches(0.3),
                    v, size=11, bold=True, color=TEXT_DARK)

    right_x = Inches(6.8)
    add_textbox(slide, right_x, Inches(2.3), Inches(6), Inches(0.4),
                "关键选型理由", size=18, bold=True, color=PRIMARY)
    add_rect(slide, right_x, Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    reasons = [
        ("Spring Boot 3 + Java 17",
         "需求指定；生态成熟，自动装配开箱即用，长期 LTS 支持。"),
        ("JPA + MyBatis-Plus 双栈",
         "JPA 写实体清晰；MP 写复杂聚合 SQL 灵活。模块边界隔离。"),
        ("MySQL 8.0",
         "原生 JSON、窗口函数、CTE 表达力强；按 user_id 分区。"),
        ("MinIO（S3 兼容）",
         "一套 API 同时支持 MinIO / OSS / S3，迁移云端零改造。"),
        ("Flyway",
         "数据库结构版本化，DDL 集中管理、可回滚、可审计。"),
    ]
    for i, (t, d) in enumerate(reasons):
        y = Inches(2.85) + i * Inches(0.82)
        add_rect(slide, right_x, y + Inches(0.15), Inches(0.1), Inches(0.1), ACCENT)
        add_textbox(slide, right_x + Inches(0.2), y, Inches(5.8), Inches(0.35),
                    t, size=13, bold=True, color=PRIMARY)
        add_textbox(slide, right_x + Inches(0.2), y + Inches(0.32), Inches(5.8), Inches(0.5),
                    d, size=10, color=TEXT_GRAY, line_spacing=1.3)


# ============== 前端技术栈 ==============
def slide_tech_frontend(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "技术选型 · 前端", "Frontend Tech Stack"
    )

    categories = [
        ("框架与构建", [
            ("TypeScript", "5.4.5 · 类型化开发"),
            ("Vue 3", "3.4.27 · Composition API"),
            ("Vite", "5.2.11 · 冷启动 / HMR 极快"),
            ("Pinia", "2.1.7 · 替代 Vuex"),
            ("Vue Router", "4.3.2 · 官方路由方案"),
        ]),
        ("UI 与交互", [
            ("Element Plus", "2.7.0 · 中后台组件库"),
            ("Axios", "1.6.8 · 拦截器统一处理"),
            ("ECharts", "5.5.0 · 统计图表"),
            ("vue-i18n", "9.13.1 · 中英文国际化"),
            ("VueUse", "10.11 · Composition 工具集"),
        ]),
        ("手写与工程化", [
            ("perfect-freehand", "1.2.2 · 笔锋还原核心"),
            ("signature_pad", "5.0.4 · 触摸轨迹采集"),
            ("ESLint + Prettier", "代码规范统一"),
            ("Vitest + Playwright", "单元 + E2E 测试"),
            ("Husky + Commitlint", "Git 钩子自动化"),
        ]),
    ]
    for i, (cat, items) in enumerate(categories):
        x = Inches(0.5) + i * Inches(4.2)
        add_rounded_rect(slide, x, Inches(2.3), Inches(4.0), Inches(4.7), BG_CARD)
        add_textbox(slide, x + Inches(0.3), Inches(2.45), Inches(3.6), Inches(0.4),
                    cat, size=15, bold=True, color=PRIMARY)
        add_rect(slide, x + Inches(0.3), Inches(2.85), Inches(0.3), Inches(0.04), ACCENT)
        for j, (n, d) in enumerate(items):
            y = Inches(3.05) + j * Inches(0.7)
            add_rect(slide, x + Inches(0.3), y + Inches(0.15), Inches(0.06), Inches(0.06), ACCENT)
            add_textbox(slide, x + Inches(0.5), y, Inches(3.4), Inches(0.32),
                        n, size=12, bold=True, color=TEXT_DARK)
            add_textbox(slide, x + Inches(0.5), y + Inches(0.3), Inches(3.4), Inches(0.32),
                        d, size=10, color=TEXT_GRAY)


# ============== 系统架构 ==============
def slide_architecture(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "系统架构总览", "Overall Architecture"
    )

    add_rounded_rect(slide, Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.8), PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(2.55), Inches(12.3), Inches(0.5),
                "客户端层   ·   Vue 3 + Vite + TypeScript + Element Plus + Pinia",
                size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(3.25), Inches(12.3), Inches(0.3),
                "↓   HTTPS / JSON   ↓", size=11, color=ACCENT,
                align=PP_ALIGN.CENTER, font=EN_FONT)

    add_rounded_rect(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(1.0),
                     BG_CARD, line=PRIMARY, line_width=Pt(1.5))
    add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.4),
                "服务层   ·   Spring Boot 3 API",
                size=15, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.5),
                "Controller  →  Service  →  Repository     |     JWT 鉴权 · 全局异常 · OpenAPI 文档",
                size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.3),
                "↓", size=14, color=ACCENT, align=PP_ALIGN.CENTER, font=EN_FONT)

    inf_y = Inches(5.0)
    inf_h = Inches(1.5)
    boxes = [
        ("MySQL 8.0", "主数据库\n事务 / 索引 / 分区", PRIMARY),
        ("Redis 7.x", "缓存 / Token / 限流\nSession / 分布式锁", RGBColor(0xB9, 0x1C, 0x1C)),
        ("MinIO / Azure", "对象存储\nS3 兼容 · 文件直传", RGBColor(0x06, 0x6B, 0x4D)),
    ]
    for i, (t, d, c) in enumerate(boxes):
        x = Inches(0.5) + i * Inches(4.2)
        add_rounded_rect(slide, x, inf_y, Inches(4.0), inf_h, c)
        add_textbox(slide, x, inf_y + Inches(0.2), Inches(4.0), Inches(0.4),
                    t, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, inf_y + Inches(0.65), Inches(4.0), Inches(0.8),
                    d, size=11, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.3)

    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.3),
                "前后端分离 · 容器化编排 · 多对象存储 Provider 适配",
                size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ============== 后端分层 ==============
def slide_backend_layers(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "后端分层架构", "Backend Layered Architecture"
    )

    layers = [
        ("api", "Controller 层", "参数接收 · 统一返回 R<T> · 异常包装", PRIMARY),
        ("application", "Service 层", "业务用例编排 · 事务控制", RGBColor(0x2C, 0x52, 0x82)),
        ("domain", "Domain 层", "实体 / 值对象 / 领域服务", RGBColor(0x3D, 0x6F, 0xA0)),
        ("infrastructure", "Infrastructure 层", "持久化 · 缓存 · 对象存储 · MQ", RGBColor(0x4D, 0x8A, 0xBE)),
        ("common / config", "公共 & 配置", "统一响应 · 错误码 · Security · OpenAPI", ACCENT),
    ]
    base_y = Inches(2.3)
    h = Inches(0.85)
    for i, (pkg, name, desc, color) in enumerate(layers):
        y = base_y + i * (h + Inches(0.1))
        add_rounded_rect(slide, Inches(0.5), y, Inches(6.5), h, color)
        add_textbox(slide, Inches(0.7), y + Inches(0.1), Inches(2.5), Inches(0.3),
                    f"com.example.handwriting.{pkg}",
                    size=10, color=WHITE, font=EN_FONT)
        add_textbox(slide, Inches(0.7), y + Inches(0.35), Inches(2.5), Inches(0.4),
                    name, size=14, bold=True, color=WHITE)
        add_textbox(slide, Inches(3.4), y + Inches(0.25), Inches(3.5), Inches(0.5),
                    desc, size=11, color=WHITE, line_spacing=1.3)

    right_x = Inches(7.5)
    add_textbox(slide, right_x, Inches(2.3), Inches(5.5), Inches(0.4),
                "依赖方向", size=16, bold=True, color=PRIMARY)
    add_rect(slide, right_x, Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    add_textbox(slide, right_x, Inches(2.9), Inches(5.5), Inches(0.4),
                "api  →  application  →  domain  ←  infrastructure",
                size=13, bold=True, color=PRIMARY, font=EN_FONT)

    add_textbox(slide, right_x, Inches(3.5), Inches(5.5), Inches(0.4),
                "设计原则", size=16, bold=True, color=PRIMARY)
    add_rect(slide, right_x, Inches(3.9), Inches(0.4), Inches(0.04), ACCENT)

    principles = [
        ("单向依赖", "上层依赖下层，下层不感知上层"),
        ("接口隔离", "Controller 只依赖 Service 接口"),
        ("统一规约", "R<T> 统一返回 / BizException 全局异常"),
        ("配置外置", "全部走 application.yml + 环境变量"),
    ]
    for i, (t, d) in enumerate(principles):
        y = Inches(4.1) + i * Inches(0.65)
        add_rect(slide, right_x, y + Inches(0.15), Inches(0.1), Inches(0.1), ACCENT)
        add_textbox(slide, right_x + Inches(0.25), y, Inches(5), Inches(0.3),
                    t, size=12, bold=True, color=PRIMARY)
        add_textbox(slide, right_x + Inches(0.25), y + Inches(0.28), Inches(5), Inches(0.3),
                    d, size=10, color=TEXT_GRAY)


# ============== 前端模块化 ==============
def slide_frontend_modules(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "前端模块化设计", "Frontend Modular Structure"
    )

    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12), Inches(0.4),
                "数据流：View  →  Pinia Action  →  API (Axios)  →  后端",
                size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    modules = [
        ("api/", "接口层", "按业务域聚合\naxios 请求封装", PRIMARY),
        ("stores/", "状态层", "Pinia 状态管理\nuser/sample/audit/dict/app", RGBColor(0x2C, 0x52, 0x82)),
        ("views/", "视图层", "portal / auth\nsample / admin", RGBColor(0x3D, 0x6F, 0xA0)),
        ("composables/", "组合式", "useAuth / useHandwriting\nuseUpload / usePermission", RGBColor(0x4D, 0x8A, 0xBE)),
    ]
    for i, (path, name, desc, color) in enumerate(modules):
        x = Inches(0.5) + i * Inches(3.1)
        add_rounded_rect(slide, x, Inches(3.1), Inches(2.9), Inches(2.0), color)
        add_textbox(slide, x + Inches(0.2), Inches(3.2), Inches(2.6), Inches(0.4),
                    path, size=18, bold=True, color=WHITE, font=EN_FONT)
        add_textbox(slide, x + Inches(0.2), Inches(3.6), Inches(2.6), Inches(0.4),
                    name, size=14, bold=True, color=ACCENT)
        add_textbox(slide, x + Inches(0.2), Inches(4.1), Inches(2.6), Inches(0.9),
                    desc, size=10, color=WHITE, line_spacing=1.4)

    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12), Inches(0.4),
                "辅助模块", size=14, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.5), Inches(5.8), Inches(0.4), Inches(0.04), ACCENT)

    aux = [
        ("components/", "通用组件 (base/business)"),
        ("layouts/", "布局 (Default/Admin/Blank)"),
        ("router/", "Vue Router + 守卫"),
        ("directives/", "自定义指令 (v-permission)"),
        ("utils/", "request / storage / format"),
        ("styles/", "SCSS 变量 + 墨韵主题"),
        ("locales/", "i18n 中英文文案"),
        ("types/", "TypeScript 类型声明"),
    ]
    for i, (path, desc) in enumerate(aux):
        col = i % 4
        row = i // 4
        x = Inches(0.5) + col * Inches(3.1)
        y = Inches(6.0) + row * Inches(0.55)
        add_rounded_rect(slide, x, y, Inches(2.9), Inches(0.45), BG_CARD)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.05), Inches(1.2), Inches(0.35),
                    path, size=11, bold=True, color=PRIMARY, font=EN_FONT)
        add_textbox(slide, x + Inches(1.4), y + Inches(0.08), Inches(1.45), Inches(0.35),
                    desc, size=9, color=TEXT_GRAY)


# ============== JWT + RBAC ==============
def slide_key_jwt(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "关键技术 ① ", "JWT 无状态鉴权 + 双 Token 机制"
    )

    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(0.4),
                "双 Token 机制", size=16, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.5), Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    add_rounded_rect(slide, Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.4), BG_CARD)
    add_textbox(slide, Inches(0.7), Inches(3.0), Inches(2), Inches(0.4),
                "Access Token", size=14, bold=True, color=PRIMARY, font=EN_FONT)
    add_textbox(slide, Inches(2.8), Inches(3.05), Inches(3.5), Inches(0.4),
                "TTL：1 小时", size=11, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.7), Inches(3.45), Inches(5.6), Inches(0.8),
                "· 承载用户身份与角色信息\n"
                "· 通过 Authorization: Bearer 头传递\n"
                "· 服务端无状态校验，无需 Session 共享",
                size=10, color=TEXT_GRAY, line_spacing=1.4)

    add_rounded_rect(slide, Inches(0.5), Inches(4.5), Inches(6.0), Inches(1.4), BG_CARD)
    add_textbox(slide, Inches(0.7), Inches(4.6), Inches(2.5), Inches(0.4),
                "Refresh Token", size=14, bold=True, color=PRIMARY, font=EN_FONT)
    add_textbox(slide, Inches(2.8), Inches(4.65), Inches(3.5), Inches(0.4),
                "TTL：7 天 · 一次性使用", size=11, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.7), Inches(5.05), Inches(5.6), Inches(0.8),
                "· 一次性使用，使用后立即轮换\n"
                "· 登出 / 改密时写入 Redis 黑名单\n"
                "· 实现「无状态 + 可吊销」双重目标",
                size=10, color=TEXT_GRAY, line_spacing=1.4)

    add_textbox(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(0.4),
                "RBAC 权限模型", size=16, bold=True, color=PRIMARY)
    add_rect(slide, Inches(6.8), Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    roles = [
        ("USER", "普通用户", "书写上传 · 查看样本 · 个人中心", RGBColor(0x05, 0x96, 0x69)),
        ("AUDITOR", "审核员", "审核工作流 · 审计日志", RGBColor(0xD9, 0x73, 0x06)),
        ("ADMIN", "管理员", "用户管理 · 字典管理 · 系统配置", RGBColor(0xB9, 0x1C, 0x1C)),
    ]
    for i, (code, name, desc, c) in enumerate(roles):
        y = Inches(2.9) + i * Inches(1.0)
        add_rounded_rect(slide, Inches(6.8), y, Inches(6.0), Inches(0.85), c)
        add_textbox(slide, Inches(7.0), y + Inches(0.15), Inches(1.5), Inches(0.4),
                    code, size=14, bold=True, color=WHITE, font=EN_FONT)
        add_textbox(slide, Inches(8.4), y + Inches(0.1), Inches(4), Inches(0.4),
                    name, size=14, bold=True, color=WHITE)
        add_textbox(slide, Inches(7.0), y + Inches(0.45), Inches(5.6), Inches(0.4),
                    desc, size=10, color=GOLD_SOFT)

    add_textbox(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
                "实现亮点：相比传统 Session 方案，JWT + Refresh 让后端真正无状态，水平扩展友好；通过 Redis 黑名单兼顾安全性。",
                size=12, color=PRIMARY, align=PP_ALIGN.CENTER, bold=True)
    add_rect(slide, Inches(4.5), Inches(6.65), Inches(4.3), Inches(0.04), ACCENT)


# ============== JPA + MP ==============
def slide_key_dual(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "关键技术 ② ", "持久层双栈：JPA + MyBatis-Plus"
    )

    add_rounded_rect(slide, Inches(0.5), Inches(2.3), Inches(6.0), Inches(2.3), BG_CARD)
    add_rect(slide, Inches(0.5), Inches(2.3), Inches(0.08), Inches(2.3), PRIMARY)
    add_textbox(slide, Inches(0.8), Inches(2.45), Inches(2), Inches(0.4),
                "JPA", size=20, bold=True, color=PRIMARY, font=EN_FONT)
    add_textbox(slide, Inches(2.0), Inches(2.55), Inches(4), Inches(0.3),
                "Spring Data JPA · 单表 CRUD", size=10, color=TEXT_GRAY)
    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(5.6), Inches(1.5),
                "✓  Repository 模式，方法命名即查询\n"
                "✓  实体建模清晰，与数据库结构自动映射\n"
                "✓  启动时 ddl-auto=validate 校验表结构\n"
                "✓  适用：User / Sample / CharDict 等标准实体",
                size=12, color=TEXT_DARK, line_spacing=1.5)

    add_rounded_rect(slide, Inches(6.8), Inches(2.3), Inches(6.0), Inches(2.3), BG_CARD)
    add_rect(slide, Inches(6.8), Inches(2.3), Inches(0.08), Inches(2.3), ACCENT)
    add_textbox(slide, Inches(7.1), Inches(2.45), Inches(3.5), Inches(0.4),
                "MyBatis-Plus", size=20, bold=True, color=ACCENT, font=EN_FONT)
    add_textbox(slide, Inches(7.1), Inches(2.95), Inches(5.6), Inches(0.3),
                "3.5.5 · 复杂 SQL 与聚合查询", size=10, color=TEXT_GRAY)
    add_textbox(slide, Inches(7.1), Inches(3.3), Inches(5.6), Inches(1.3),
                "✓  Lambda Wrapper 灵活拼装 SQL\n"
                "✓  支持复杂联表、聚合、分组、统计\n"
                "✓  内置分页插件、逻辑删除、审计字段\n"
                "✓  适用：Dashboard 统计 / 趋势 / TOP N",
                size=12, color=TEXT_DARK, line_spacing=1.5)

    add_rect(slide, Inches(6.55), Inches(2.4), Inches(0.05), Inches(2.1), DIVIDER)

    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
                "共存策略：模块边界隔离，绝不混用", size=16, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.5), Inches(5.4), Inches(0.4), Inches(0.04), ACCENT)

    rules = [
        ("实体类用 JPA", "@Entity 标注，由 JPA 管理"),
        ("自定义查询走 MP", "Mapper + Lambda Wrapper"),
        ("事务统一管理", "@Transactional 加在 Service 层"),
        ("绝不交叉依赖", "Repository 不调用 Mapper，反之亦然"),
    ]
    for i, (t, d) in enumerate(rules):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(5.65) + row * Inches(0.7)
        add_rounded_rect(slide, x, y, Inches(6.0), Inches(0.6), BG_CARD)
        add_rect(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.1), Inches(0.2), ACCENT)
        add_textbox(slide, x + Inches(0.4), y + Inches(0.05), Inches(2.5), Inches(0.5),
                    t, size=13, bold=True, color=PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + Inches(2.9), y + Inches(0.05), Inches(3), Inches(0.5),
                    d, size=11, color=TEXT_GRAY, anchor=MSO_ANCHOR.MIDDLE)


# ============== 文件直传 ==============
def slide_key_storage(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "关键技术 ③ ", "文件直传 + 多对象存储 Provider"
    )

    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(0.4),
                "文件直传流程", size=16, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.5), Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    flow = [
        ("1", "客户端请求", "GET /v1/file/sign\n获取 OSS 预签名 URL"),
        ("2", "客户端直传", "PUT 文件到 MinIO/S3\n不经过后端中转"),
        ("3", "元数据提交", "POST /v1/sample/upload\nfileKey · charId · sha256"),
        ("4", "服务端校验", "HeadObject 校验文件\n落库 + MQ 异步任务"),
    ]
    for i, (no, t, d) in enumerate(flow):
        y = Inches(2.9) + i * Inches(0.95)
        add_rect(slide, Inches(0.5), y + Inches(0.1), Inches(0.5), Inches(0.5), ACCENT)
        add_textbox(slide, Inches(0.5), y + Inches(0.15), Inches(0.5), Inches(0.4),
                    no, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=EN_FONT)
        add_textbox(slide, Inches(1.2), y + Inches(0.05), Inches(5.0), Inches(0.35),
                    t, size=13, bold=True, color=PRIMARY)
        add_textbox(slide, Inches(1.2), y + Inches(0.4), Inches(5.0), Inches(0.5),
                    d, size=10, color=TEXT_GRAY, line_spacing=1.3)
        if i < 3:
            add_rect(slide, Inches(0.72), y + Inches(0.6), Inches(0.06), Inches(0.3), DIVIDER)

    add_textbox(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(0.4),
                "多 Provider 适配", size=16, bold=True, color=PRIMARY)
    add_rect(slide, Inches(6.8), Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    providers = [
        ("Local", "本地文件系统（开发）", RGBColor(0x6B, 0x72, 0x80)),
        ("MinIO", "S3 兼容（自建）", RGBColor(0xC7, 0x2E, 0x2E)),
        ("S3", "AWS S3（云端）", RGBColor(0xD9, 0x73, 0x06)),
        ("Azure", "Azure Blob Storage", RGBColor(0x05, 0x6B, 0xA0)),
    ]
    for i, (p, d, c) in enumerate(providers):
        col = i % 2
        row = i // 2
        x = Inches(6.8) + col * Inches(3.05)
        y = Inches(2.9) + row * Inches(0.95)
        add_rounded_rect(slide, x, y, Inches(2.9), Inches(0.85), c)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.1), Inches(2.5), Inches(0.4),
                    p, size=15, bold=True, color=WHITE, font=EN_FONT)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.45), Inches(2.5), Inches(0.4),
                    d, size=10, color=GOLD_SOFT)

    add_textbox(slide, Inches(6.8), Inches(4.9), Inches(6), Inches(0.4),
                "分桶策略", size=14, bold=True, color=PRIMARY)
    add_rect(slide, Inches(6.8), Inches(5.3), Inches(0.4), Inches(0.04), ACCENT)

    buckets = [
        ("handwriting-collect", "采集桶", "用户书写样本原图"),
        ("handwriting-audit", "审核桶", "审核员工作流附件"),
        ("handwriting-backup", "备份桶", "定期备份与归档"),
    ]
    for i, (b, t, d) in enumerate(buckets):
        y = Inches(5.5) + i * Inches(0.4)
        add_textbox(slide, Inches(6.8), y, Inches(2.5), Inches(0.35),
                    b, size=11, bold=True, color=PRIMARY, font=EN_FONT)
        add_textbox(slide, Inches(9.4), y, Inches(0.8), Inches(0.35),
                    t, size=11, color=ACCENT, bold=True)
        add_textbox(slide, Inches(10.3), y, Inches(2.5), Inches(0.35),
                    d, size=10, color=TEXT_GRAY)

    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                "优势：解放后端带宽 · 降低服务器 IO · 支持大文件 · 断点续传友好 · 一套配置切换存储后端",
                size=12, color=PRIMARY, align=PP_ALIGN.CENTER, bold=True)


# ============== 手写板 ==============
def slide_key_handwriting(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "关键技术 ④ ", "手写板笔锋还原（核心亮点）"
    )

    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(0.4),
                "核心技术栈", size=16, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.5), Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    techs = [
        ("perfect-freehand", "变半径贝塞尔拟合算法",
         "基于压力 + 速度还原毛笔 / 钢笔笔锋"),
        ("signature_pad", "PointerEvent 三维数据采集",
         "x / y / pressure 实时采集"),
        ("Canvas DPR 适配", "devicePixelRatio 适配",
         "高分辨率屏不模糊"),
        ("SVG 序列化", "笔画先序列化再栅格化",
         "笔顺可还原、可二次编辑"),
    ]
    for i, (t, st, d) in enumerate(techs):
        y = Inches(2.9) + i * Inches(0.85)
        add_rounded_rect(slide, Inches(0.5), y, Inches(6.0), Inches(0.75), BG_CARD)
        add_rect(slide, Inches(0.5), y, Inches(0.08), Inches(0.75), ACCENT)
        add_textbox(slide, Inches(0.7), y + Inches(0.05), Inches(2.5), Inches(0.32),
                    t, size=12, bold=True, color=PRIMARY, font=EN_FONT)
        add_textbox(slide, Inches(0.7), y + Inches(0.32), Inches(2.5), Inches(0.4),
                    st, size=10, color=ACCENT)
        add_textbox(slide, Inches(3.3), y + Inches(0.15), Inches(3.1), Inches(0.5),
                    d, size=10, color=TEXT_GRAY, line_spacing=1.3)

    add_textbox(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(0.4),
                "功能特性", size=16, bold=True, color=PRIMARY)
    add_rect(slide, Inches(6.8), Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    features = [
        ("笔刷预设", "钢笔 / 毛笔 / 铅笔\n三套参数实时切换"),
        ("撤销 / 重做", "基于命令模式\n操作栈管理"),
        ("压感事件", "支持 Apple Pencil\nWacom 等真实压感"),
        ("高清导出", "SVG 矢量 + PNG 位图\n双格式输出"),
    ]
    for i, (t, d) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(6.8) + col * Inches(3.05)
        y = Inches(2.9) + row * Inches(1.4)
        add_rounded_rect(slide, x, y, Inches(2.9), Inches(1.3), PRIMARY)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.4),
                    t, size=14, bold=True, color=ACCENT)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.5), Inches(0.7),
                    d, size=10, color=WHITE, line_spacing=1.3)

    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                "算法原理：PointerEvent.pressure  →  变半径拟合  →  贝塞尔曲线光栅化  →  Canvas / SVG",
                size=12, color=PRIMARY, align=PP_ALIGN.CENTER, bold=True)
    add_rect(slide, Inches(2.0), Inches(6.85), Inches(9.3), Inches(0.03), ACCENT)


# ============== Flyway + 安全 ==============
def slide_key_flyway_security(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "关键技术 ⑤⑥ ", "Flyway 数据库迁移 & 安全防护"
    )

    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(6), Inches(0.4),
                "Flyway 数据库版本化迁移", size=15, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.5), Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    scripts = [
        ("V1__init_schema.sql", "初始化核心表结构"),
        ("V2__seed_data.sql", "种子数据"),
        ("V3__widen_tinyint_to_int.sql", "字段类型扩展"),
        ("V4__password_recovery.sql", "密码恢复相关表"),
        ("V5__align_recovery_column_types.sql", "字段对齐"),
    ]
    for i, (n, d) in enumerate(scripts):
        y = Inches(2.9) + i * Inches(0.55)
        add_rounded_rect(slide, Inches(0.5), y, Inches(6.0), Inches(0.45), BG_CARD)
        add_rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.45), ACCENT)
        add_textbox(slide, Inches(0.7), y + Inches(0.05), Inches(2.8), Inches(0.35),
                    n, size=10, bold=True, color=PRIMARY, font=EN_FONT)
        add_textbox(slide, Inches(3.5), y + Inches(0.08), Inches(2.9), Inches(0.35),
                    d, size=10, color=TEXT_GRAY)

    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(6), Inches(0.4),
                "解决：开发 / 测试 / 生产环境数据库结构不一致",
                size=11, color=PRIMARY, bold=True)
    add_textbox(slide, Inches(0.5), Inches(6.2), Inches(6), Inches(0.4),
                "· 启动时按版本号顺序执行   · DDL 集中管理、可回滚、可审计",
                size=10, color=TEXT_GRAY)

    add_textbox(slide, Inches(6.8), Inches(2.3), Inches(6), Inches(0.4),
                "安全防护体系", size=15, bold=True, color=PRIMARY)
    add_rect(slide, Inches(6.8), Inches(2.7), Inches(0.4), Inches(0.04), ACCENT)

    security = [
        ("图形验证码", "easy-captcha · Redis · TTL 5min", RGBColor(0x05, 0x96, 0x69)),
        ("接口限流", "Redis 令牌桶 · 20/min · 60/min", RGBColor(0xD9, 0x73, 0x06)),
        ("TOTP 二次验证", "ZXing 二维码 · RFC 6238", RGBColor(0x05, 0x6B, 0xA0)),
        ("密保问题", "用户枚举防护 · 时间归一化", RGBColor(0xB9, 0x1C, 0x1C)),
        ("BCrypt 密码哈希", "strength=10 · 加盐存储", PRIMARY),
    ]
    for i, (t, d, c) in enumerate(security):
        y = Inches(2.9) + i * Inches(0.55)
        add_rounded_rect(slide, Inches(6.8), y, Inches(6.0), Inches(0.45), c)
        add_textbox(slide, Inches(7.0), y + Inches(0.05), Inches(2), Inches(0.35),
                    t, size=11, bold=True, color=WHITE)
        add_textbox(slide, Inches(9.0), y + Inches(0.08), Inches(3.7), Inches(0.35),
                    d, size=9, color=GOLD_SOFT)

    add_textbox(slide, Inches(6.8), Inches(5.8), Inches(6), Inches(0.4),
                "纵深防御：密码强度 + 验证码 + 限流 + 二次验证",
                size=11, color=PRIMARY, bold=True)
    add_textbox(slide, Inches(6.8), Inches(6.2), Inches(6), Inches(0.4),
                "· 抵御爆破 · 防用户枚举 · 关键操作二次确认",
                size=10, color=TEXT_GRAY)


# ============== 工程化 ==============
def slide_engineering(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "工程化与部署", "Engineering & Deployment"
    )

    items = [
        ("统一规约", "R<T> 统一返回 · BizException 全局异常\nErrorCode 错误码枚举 · CORS 配置",
         PRIMARY),
        ("OpenAPI 契约", "springdoc-openapi 自动生成\nSwagger UI 在线调试\n前后端契约式开发",
         RGBColor(0x2C, 0x52, 0x82)),
        ("Actuator 监控", "/actuator/health · /actuator/prometheus\nJVM / HTTP / DB 指标\n可对接 Prometheus + Grafana",
         RGBColor(0x3D, 0x6F, 0xA0)),
        ("代码规范", "ESLint + Prettier + Stylelint\nHusky + lint-staged 卡点\nCommitlint 强制规范",
         RGBColor(0x4D, 0x8A, 0xBE)),
    ]
    for i, (t, d, c) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(2.3) + row * Inches(2.0)
        add_rounded_rect(slide, x, y, Inches(6.0), Inches(1.85), c)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.5), Inches(0.4),
                    t, size=18, bold=True, color=WHITE)
        add_rect(slide, x + Inches(0.3), y + Inches(0.7), Inches(0.3), Inches(0.04), ACCENT)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.85), Inches(5.5), Inches(1.0),
                    d, size=12, color=GOLD_SOFT, line_spacing=1.5)

    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                "容器化部署：Dockerfile 多阶段构建  ·  docker-compose 一键起 MySQL/Redis/MinIO/Client  ·  K8s Deployment YAML",
                size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER, line_spacing=1.4)
    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
                "docker compose -f deploy/docker-compose.yml --env-file .env up -d",
                size=10, color=ACCENT, align=PP_ALIGN.CENTER, font=EN_FONT)


# ============== 测试 ==============
def slide_testing(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "测试与质量保证", "Testing & Quality Assurance"
    )

    cols = [
        ("后端测试", "JUnit 5 + Mockito", [
            "CharDict 解析单元测试",
            "批量导入业务测试",
            "Spring Security 集成测试",
            "Repository 边界用例覆盖",
        ], PRIMARY),
        ("前端测试", "Vitest + jsdom", [
            "composables 单元测试",
            "useHandwriting 手写逻辑",
            "utils 工具函数全覆盖",
            "storage / validator / format",
        ], RGBColor(0x2C, 0x52, 0x82)),
        ("E2E & 联调", "Playwright + Postman", [
            "关键流程回归（登录/上传/审核）",
            "跨浏览器 / 多端适配",
            "Postman 收录 60+ 用例",
            "全链路接口联调",
        ], RGBColor(0x3D, 0x6F, 0xA0)),
    ]
    for i, (t, sub, items, c) in enumerate(cols):
        x = Inches(0.5) + i * Inches(4.2)
        add_rounded_rect(slide, x, Inches(2.3), Inches(4.0), Inches(4.7), BG_CARD)
        add_rect(slide, x, Inches(2.3), Inches(4.0), Inches(0.1), c)
        add_textbox(slide, x + Inches(0.3), Inches(2.5), Inches(3.6), Inches(0.5),
                    t, size=20, bold=True, color=PRIMARY)
        add_textbox(slide, x + Inches(0.3), Inches(3.0), Inches(3.6), Inches(0.3),
                    sub, size=11, color=ACCENT, bold=True)
        add_rect(slide, x + Inches(0.3), Inches(3.4), Inches(0.4), Inches(0.04), ACCENT)
        for j, item in enumerate(items):
            y = Inches(3.65) + j * Inches(0.7)
            add_rect(slide, x + Inches(0.3), y + Inches(0.1), Inches(0.1), Inches(0.1), ACCENT)
            add_textbox(slide, x + Inches(0.5), y, Inches(3.4), Inches(0.5),
                        item, size=12, color=TEXT_DARK, line_spacing=1.3)

    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
                "质量门禁：自动化测试 + Code Review + 提交前 Lint，确保核心模块稳定性",
                size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ============== 亮点 ==============
def slide_highlights(prs, page_no, total_pages, section_no, section_title):
    slide = slide_content_base(
        prs, page_no, total_pages, section_no, section_title,
        "亮点与难点总结", "Highlights & Key Challenges"
    )

    highlights = [
        ("01", "JPA + MP 双栈共存", "通过模块边界隔离，兼顾开发效率与复杂查询表达力", PRIMARY),
        ("02", "手写体笔锋还原", "perfect-freehand 变半径拟合，真实还原书写压感", RGBColor(0x2C, 0x52, 0x82)),
        ("03", "多 Provider 对象存储", "MinIO / S3 / Azure / 本地 4 套实现，一套配置切换", RGBColor(0x3D, 0x6F, 0xA0)),
        ("04", "完备的安全体系", "JWT + Refresh + Redis 黑名单 + BCrypt + 验证码 + TOTP", RGBColor(0x4D, 0x8A, 0xBE)),
        ("05", "工程化规范落地", "Flyway · OpenAPI 契约 · ESLint/Stylelint · Husky 自动化", RGBColor(0x05, 0x96, 0x69)),
        ("06", "可观测性", "Actuator + Prometheus · JWT 审计 · Spring Security Debug", RGBColor(0xD9, 0x73, 0x06)),
    ]
    for i, (no, t, d, c) in enumerate(highlights):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(2.3) + row * Inches(1.55)
        add_rounded_rect(slide, x, y, Inches(6.0), Inches(1.4), c)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(1), Inches(0.8),
                    no, size=40, bold=True, color=ACCENT, font=EN_FONT)
        add_textbox(slide, x + Inches(1.3), y + Inches(0.2), Inches(4.5), Inches(0.4),
                    t, size=16, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(1.3), y + Inches(0.65), Inches(4.5), Inches(0.7),
                    d, size=11, color=GOLD_SOFT, line_spacing=1.4)

    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                "在「前后端分离 · 持久层双栈 · 对象存储适配 · 安全防护 · 工程化规范」方面进行了完整的企业级工程实践",
                size=12, color=PRIMARY, align=PP_ALIGN.CENTER, bold=True)


# ============== Q&A ==============
def slide_qa(prs, page_no, total_pages):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)

    add_rect(slide, 0, Inches(3.4), SLIDE_W, Inches(0.04), ACCENT)
    add_rect(slide, 0, Inches(4.2), SLIDE_W, Inches(0.02), ACCENT)

    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
                "Q & A", size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=EN_FONT)
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.6),
                "欢迎各位老师批评指正", size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.4),
                "Thank you for your attention and guidance.",
                size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, font=EN_FONT)

    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.5),
                "特别致谢", size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.4),
                "指导教师 XXX 教授   ·   课程组全体老师   ·   项目组同学",
                size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
                "开源社区：Spring / Vue / Element Plus / perfect-freehand",
                size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
                f"HandWrite System  |  {page_no} / {total_pages}",
                size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, font=EN_FONT)


# ============== 主流程 ==============
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 21 页 PPT
    TOTAL = 21

    # 1. 封面
    slide_cover(prs)
    # 2. 目录
    slide_toc(prs, 2, TOTAL)
    # 3. 章节 01 - 项目背景
    slide_section(prs, 3, TOTAL, 1, "项目背景与目标", "Background & Goals")
    # 4. 项目背景详情
    slide_background(prs, 4, TOTAL, 1, "项目背景与目标")
    # 5. 章节 02 - 技术选型
    slide_section(prs, 5, TOTAL, 2, "技术选型", "Technology Stack")
    # 6. 后端技术栈
    slide_tech_backend(prs, 6, TOTAL, 2, "技术选型")
    # 7. 前端技术栈
    slide_tech_frontend(prs, 7, TOTAL, 2, "技术选型")
    # 8. 章节 03 - 系统架构
    slide_section(prs, 8, TOTAL, 3, "系统架构设计", "System Architecture")
    # 9. 架构总览
    slide_architecture(prs, 9, TOTAL, 3, "系统架构")
    # 10. 后端分层
    slide_backend_layers(prs, 10, TOTAL, 3, "系统架构")
    # 11. 前端模块
    slide_frontend_modules(prs, 11, TOTAL, 3, "系统架构")
    # 12. 章节 04 - 关键技术
    slide_section(prs, 12, TOTAL, 4, "关键技术实现", "Key Implementation")
    # 13. JWT + RBAC
    slide_key_jwt(prs, 13, TOTAL, 4, "关键技术")
    # 14. JPA + MP
    slide_key_dual(prs, 14, TOTAL, 4, "关键技术")
    # 15. 文件直传
    slide_key_storage(prs, 15, TOTAL, 4, "关键技术")
    # 16. 手写板
    slide_key_handwriting(prs, 16, TOTAL, 4, "关键技术")
    # 17. Flyway + 安全
    slide_key_flyway_security(prs, 17, TOTAL, 4, "关键技术")
    # 18. 工程化与部署
    slide_engineering(prs, 18, TOTAL, 5, "工程化与部署")
    # 19. 测试与质量
    slide_testing(prs, 19, TOTAL, 6, "测试与质量保证")
    # 20. 亮点总结
    slide_highlights(prs, 20, TOTAL, 7, "亮点与展望")
    # 21. Q&A
    slide_qa(prs, 21, TOTAL)

    out_path = r"d:\dev\JavaCourseDesign\答辩PPT-技术篇.pptx"
    prs.save(out_path)
    print(f"[OK] PPT 已生成：{out_path}")
    print(f"[OK] 总页数：{len(prs.slides)}")


if __name__ == "__main__":
    main()
